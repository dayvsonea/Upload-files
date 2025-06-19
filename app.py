import os, tempfile, subprocess, datetime, uuid, shutil, io, re
from flask import Flask, render_template, request, send_from_directory, redirect, url_for, flash
from pptx import Presentation
import pandas as pd, requests

app = Flask(__name__)
app.secret_key = "ea-secret"

TEMPLATE_DIR = os.getenv("PPT_TEMPLATE_DIR", "./templates_ppt")
WORKBOOK_XLSX = os.getenv("WORKBOOK_XLSX", "")
UPLOAD_LIMIT = int(os.getenv("UPLOAD_LIMIT", "5"))

OPTION_MAP = {
    "1": "0000 1 Nova Venda COM Arquivo.pptx",
    "2": "0000 2 Nova Venda SEM Arquivo.pptx",
    "3": "0000 3 Upgrade COM Arquivo.pptx",
    "4": "0000 4 Upgrade SEM Arquivo.pptx"
}

REGEX = re.compile(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}")

def load_workbook():
    if not WORKBOOK_XLSX:
        return pd.DataFrame(), pd.DataFrame()
    data = requests.get(WORKBOOK_XLSX, timeout=20).content
    xls = pd.ExcelFile(io.BytesIO(data))
    planos = pd.read_excel(xls, sheet_name="Planos")
    variaveis = pd.read_excel(xls, sheet_name="Variáveis")
    return planos, variaveis

planos_df, variaveis_df = load_workbook()

def get_down_labels():
    if variaveis_df.empty:
        return [f"Arquivo {i}" for i in range(1,UPLOAD_LIMIT+1)]
    subset = variaveis_df[variaveis_df["Variável"].str.startswith("down")]
    labels = subset["Qual conteúdo"].tolist()
    return labels + [f"Arquivo {i}" for i in range(len(labels)+1, UPLOAD_LIMIT+1)]

DOWN_LABELS = get_down_labels()

def substitute_ppt(tpl, mapping, out_path):
    prs = Presentation(tpl)
    def dist(par,new):
        idx=0
        for run in par.runs:
            ln=len(run.text); run.text=new[idx:idx+ln]; idx+=ln
        if idx<len(new): par.runs[-1].text+=new[idx:]

    def proc(par, sh):
        txt="".join(r.text for r in par.runs)
        if "{{" in txt:
            def repl(m):
                k=m.group(1).lower()
                return "Baixar Arquivo Aqui" if k.startswith("down") else mapping.get(k,m.group(0))
            dist(par, REGEX.sub(repl, txt))
            m_dl=REGEX.fullmatch(txt.strip())
            if m_dl and m_dl.group(1).lower().startswith("down"):
                sh.click_action.hyperlink.address = mapping[m_dl.group(1).lower()]
        if mapping.get("valorsemdesc")==mapping.get("valorcomdesc"):
            low=txt.lower()
            if "de:" in low: [setattr(r,"text","") for r in par.runs]
            elif "por:" in low: dist(par,mapping["valorsemdesc"])

    def walk(s):
        if s.has_text_frame:
            for p in s.text_frame.paragraphs: proc(p,s)
        if s.has_table:
            for row in s.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs: proc(p,s)
        if s.shape_type==6:
            for sub in s.shapes: walk(sub)
    for slide in prs.slides:
        for shp in slide.shapes:
            walk(shp)
    prs.save(out_path)

def convert_pdf(ppt_path):
    pdf_path = ppt_path.replace(".pptx",".pdf")
    try:
        subprocess.run(["soffice","--headless","--convert-to","pdf",ppt_path,"--outdir",os.path.dirname(ppt_path)],check=True,timeout=90)
        return pdf_path
    except Exception: return ""

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/form/<opt>")
def form(opt):
    if opt not in OPTION_MAP: return redirect(url_for("index"))
    show_uploads = opt in ["1","3"]
    return render_template("form.html", opt=opt, template_name=OPTION_MAP[opt], show_uploads=show_uploads, down_labels=DOWN_LABELS)

@app.route("/generate/<opt>", methods=["POST"])
def generate(opt):
    if opt not in OPTION_MAP: return redirect(url_for("index"))
    m = {k.lower():v for k,v in request.form.items()}
    if 'desconto' not in request.form: m["valorcomdesc"]=m.get("valorsemdesc")
    # csv plan data
    row = planos_df[planos_df["Plano"].str.strip().str.lower()==m.get("plano","").lower()]
    if not row.empty:
        r=row.iloc[0]
        m.setdefault("descplan", r.get("Descrição",""))
        m.setdefault("dimensionamento", r.get("Dimensionamento",""))
    m["dataenvio"]=datetime.datetime.now().strftime("%d/%m/%Y")
    m["hoje"]=m["dataenvio"]
    for i in range(1,UPLOAD_LIMIT+1):
        m[f"down{i}"]="https://www.e-auditoria.com.br"
    tmp=tempfile.mkdtemp()
    tpl=os.path.join(TEMPLATE_DIR, OPTION_MAP[opt])
    out_ppt=os.path.join(tmp, f"Proposta_{uuid.uuid4().hex}.pptx")
    substitute_ppt(tpl,m,out_ppt)
    pdf=convert_pdf(out_ppt)
    email_txt=render_template("email_template.txt", **m)
    email_path=os.path.join(tmp,"mensagem_email.txt")
    open(email_path,"w",encoding="utf-8").write(email_txt)
    files=[("PPTX",os.path.basename(out_ppt))]
    if pdf: files.append(("PDF",os.path.basename(pdf)))
    files.append(("Mensagem (.txt)",os.path.basename(email_path)))
    return render_template("success.html", files=files, folder=os.path.basename(tmp))

@app.route("/download/<folder>/<filename>")
def download(folder, filename):
    path=os.path.join(tempfile.gettempdir(), folder)
    return send_from_directory(path, filename, as_attachment=True)

if __name__=="__main__":
    app.run(host="0.0.0.0", port=5000)
